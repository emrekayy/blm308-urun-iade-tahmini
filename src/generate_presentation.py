"""BLM308 final sunumunu oluşturur."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from src.utils import DATA_PROCESSED_DIR, FIGURES_DIR  # noqa: E402

STUDENT_NAME = "Emre Kaya"
STUDENT_NUMBER = "231041045"


def load_json(path: Path) -> dict:
    with path.open(encoding="utf-8") as file:
        return json.load(file)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.3), width=Inches(8.0))
    if caption:
        textbox = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(8.0), Inches(0.5))
        textbox.text_frame.text = caption


def build_presentation() -> Path:
    preprocessing = load_json(DATA_PROCESSED_DIR / "preprocessing_metadata.json")
    evaluation = load_json(DATA_PROCESSED_DIR / "evaluation_results.json")

    inspection = preprocessing["inspection"]
    class_dist = inspection["class_distribution"]
    comparison = evaluation["comparison_table"]
    best_model = evaluation["best_model"]
    test_metrics = evaluation["test_metrics"]
    cm = evaluation["confusion_matrix"]

    top3 = comparison[:3]
    comparison_lines = [
        f"{row['Model']}: Doğruluk={row['Accuracy']:.3f}, F1={row['F1-Score']:.3f}, AUC={row['ROC-AUC']:.3f}"
        for row in top3
    ]

    prs = Presentation()

    add_title_slide(
        prs,
        "Çevrimiçi Alışveriş Ürün İade Tahmini",
        "BLM308 Veri Madenciliği Final Projesi\nEmre Kaya | 231041045 | Mayıs 2026",
    )

    add_bullet_slide(
        prs,
        "Problem Tanımı",
        [
            "Bir e-ticaret siparişinin/ürünün iade edilip edilmeyeceğini tahmin etmek.",
            "İkili hedef: return_status (0 = iade yok, 1 = iade var).",
            "Amaç: iade kaynaklı operasyonel maliyeti azaltmak ve müşteri deneyimini iyileştirmek.",
        ],
    )

    add_bullet_slide(
        prs,
        "Motivasyon",
        [
            "İadeler lojistik ve yeniden stoklama maliyetlerini artırır.",
            "Erken risk tespiti proaktif müşteri hizmetini destekler.",
            "Elde edilen içgörüler fiyatlandırma ve ürün kalitesi iyileştirmelerine rehberlik eder.",
        ],
    )

    add_bullet_slide(
        prs,
        "Veri Seti",
        [
            f"Dosya: {preprocessing['dataset_file']}",
            f"Kayıt sayısı: {inspection['shape'][0]} | Sütunlar: {', '.join(inspection['columns'])}",
            f"Hedef değişken: {inspection['target_column']}",
            f"Sınıf dağılımı: 0={class_dist.get(0, class_dist.get('0'))}, "
            f"1={class_dist.get(1, class_dist.get('1'))}",
            "Eksik değer: tespit edilmedi.",
        ],
    )

    add_image_slide(
        prs,
        "EDA Bulguları",
        FIGURES_DIR / "01_class_distribution.png",
        "Dengesiz veri seti: yaklaşık %27 iade oranı.",
    )
    add_image_slide(
        prs,
        "Fiyat ve Puan Kalıpları",
        FIGURES_DIR / "03_price_vs_return.png",
        "Düşük puanlar ve fiyat farklılıkları iade grupları arasında görülebilir.",
    )

    add_bullet_slide(
        prs,
        "Ön İşleme",
        [
            "order_id çıkarıldı; product_id ordinal özellik olarak kullanılmadı.",
            "product_id_frequency: ürünün eğitim kümesindeki görülme sayısı.",
            "product_return_rate: ürün iade oranı yalnızca eğitim kümesinden hesaplandı.",
            "Scaler yalnızca eğitim kümesinde fit edildi; split encoding öncesinde yapıldı.",
        ],
    )

    add_bullet_slide(
        prs,
        "Modeller",
        [
            "Decision Tree (Karar Ağacı)",
            "Random Forest (Rastgele Orman)",
            "Logistic Regression (Lojistik Regresyon)",
            "Naive Bayes",
            "KNN",
        ],
    )

    add_bullet_slide(
        prs,
        "Değerlendirme Kurulumu",
        [
            "10 katmanlı Stratified CV yalnızca eğitim kümesinde; ürün encoding her fold içinde yeniden hesaplandı.",
            "Test kümesi (300 kayıt) modele dahil edilmedi.",
            "Metrikler: Doğruluk, Kesinlik, Duyarlılık, F1-skoru, ROC-AUC.",
            "En iyi model CV F1-skoruna göre seçildi; test yalnızca bir kez kullanıldı.",
        ],
    )

    add_bullet_slide(prs, "Sonuçlar (En İyi 3 Model)", comparison_lines)

    add_bullet_slide(
        prs,
        f"En İyi Model: {best_model}",
        [
            f"Doğruluk: {test_metrics['accuracy']:.4f}",
            f"Kesinlik: {test_metrics['precision']:.4f}",
            f"Duyarlılık: {test_metrics['recall']:.4f}",
            f"F1-skoru: {test_metrics['f1_score']:.4f}",
            f"ROC-AUC: {test_metrics['roc_auc']:.4f}",
        ],
    )

    add_image_slide(
        prs,
        "Confusion Matrix",
        FIGURES_DIR / "08_confusion_matrix_best_model.png",
        f"TN={cm[0][0]}, FP={cm[0][1]}, FN={cm[1][0]}, TP={cm[1][1]}",
    )

    add_bullet_slide(
        prs,
        "İş Yorumu",
        [
            "Puan ve ürün bazlı kalıplar iade riskini belirlemede etkilidir.",
            "Düşük puanlı siparişler proaktif destek için işaretlenebilir.",
            "Yüksek riskli ürünler kalite incelemesi için önceliklendirilebilir.",
        ],
    )

    add_bullet_slide(
        prs,
        "Kısıtlamalar",
        [
            "Kategori, indirim ve teslimat süresi özellikleri bulunmamaktadır.",
            "Sınıf dengesizliği azınlık sınıfının duyarlılığını etkiler.",
            "product_id kodlaması yeni ürünler için genellenebilir olmayabilir.",
        ],
    )

    add_bullet_slide(
        prs,
        "Gelecek Çalışmalar",
        [
            "Yorum metni ve teslimat özelliklerinin eklenmesi.",
            "Dengesizlik yönetimi tekniklerinin uygulanması (SMOTE).",
            "Ödeme akışında gerçek zamanlı skorlama.",
        ],
    )

    add_bullet_slide(
        prs,
        "Sonuç",
        [
            f"{best_model} modeli güçlü iade tahmin performansı göstermiştir.",
            "Proje BLM308 sınıflandırma ve değerlendirme gereksinimlerini karşılamaktadır.",
            "Elde edilen içgörüler e-ticaret iade yönetimini destekler.",
        ],
    )

    output_dir = PROJECT_ROOT / "presentation"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "BLM308_Final_Presentation.pptx"
    prs.save(output_path)
    return output_path


def main() -> None:
    path = build_presentation()
    print(f"Sunum kaydedildi: {path}")


if __name__ == "__main__":
    main()
